import json
import libs.utils as utils
import collections 
import collections.abc
from pptx import Presentation
from random import randint
from pathlib import Path
import numpy as np
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create():
    master_template_path = Path(__file__).with_name('template.pptx')
    return Presentation(master_template_path)

def save(report, name):
    try:
        report.save(f'./{name}.pptx')
    except Exception as e:
        save(report, f'./{name}-{randint(2,9) * randint(2,9) * randint(2,9)}.pptx')

def get_title_slide(report):
    return report.slides[0]

def get_layouts(report):
    return report.slide_layouts

def get_layout(report, name):
    return get_layouts(report).get_by_name(name)

def update_title(the_slide, title_str):
    title = the_slide.shapes.title
    title.text = str(title_str)

# Text specific functions
def update_text_placeholder(the_slide, placeholder_id, content, size=None):
    ph = the_slide.placeholders[placeholder_id]
    ph.text = str(content)
    if size != None:
        for paragraph in ph.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size)

# Image specific functions
def update_image_placeholder(the_slide, placeholder_id, image, scale_image=False):
    ph = the_slide.placeholders[placeholder_id]
    placeholder_height = ph.height
    placeholder_width = ph.width
    placeholder_top = ph.top
    placeholder_left = ph.left

    picture = ph.insert_picture(image)
    image_width, image_height = picture.image.size

    # Calculate the aspect ratios
    placeholder_aspect_ratio = placeholder_width / placeholder_height
    image_aspect_ratio = image_width / image_height

    # Check if scaling is necessary based on aspect ratio
    if scale_image:
        if image_aspect_ratio > placeholder_aspect_ratio:
            # Calculate the scaling factor to fit the width
            scale_factor = placeholder_width / image_width
        else:
            # Calculate the scaling factor to fit the height
            scale_factor = placeholder_height / image_height

        # Calculate the new dimensions while maintaining aspect ratio
        new_width = int(image_width * scale_factor)
        new_height = int(image_height * scale_factor)

        # Set the new dimensions of the image
        picture.width = new_width
        picture.height = new_height

        # Set the position of the picture to the bottom right corner
        picture.left = placeholder_left
        picture.top = placeholder_top

    # Reset cropping
    picture.crop_top = 0
    picture.crop_left = 0
    picture.crop_bottom = 0
    picture.crop_right = 0

# Table specific functions
def update_table_placeholder(the_slide, placeholder_id, data, size=None):
    data = np.array(data)
    ph = the_slide.placeholders[placeholder_id]
    shape = ph.insert_table(rows=data.shape[0], cols=data.shape[1])
    table = shape.table
    for row in range(0, data.shape[0]):
        for col in range(0, data.shape[1]):
            tc = table.cell(row, col)
            tc.text = data[row][col]
            for paragraph in tc.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                if size != None:
                    for run in paragraph.runs:
                        run.font.size = Pt(size)

# Title specific functions
def update_title_slide(the_slide, report_title, report_subtitle):
    update_title(the_slide, report_title)
    update_text_placeholder(the_slide, 1, report_subtitle)

def update_notes_slide(the_slide, title, content):
    update_title(the_slide, title)
    update_text_placeholder(the_slide, 13, content)

    ph = the_slide.placeholders[13]
    tf = ph.text_frame
    first = True
    for note in content:
        if first:
            tf.text = str(note)
        else:
            p = tf.add_paragraph()
            p.text = str(note)
        first = False
        # p.level = 1

def update_device_report_slide(the_slide, title, content):
    update_title(the_slide, title)
    update_text_placeholder(the_slide, 13, content)

    ph = the_slide.placeholders[13]
    tf = ph.text_frame
    tf.text = content["title"]
    for point in content["points"]:
        p = tf.add_paragraph()
        p.text = str(point)
        p.level = 1

    tf.add_paragraph()
    date = tf.add_paragraph()
    date.text = "Date range"
    date.level = 0

    dateP = tf.add_paragraph()
    dateP.text = "All data in this report is from " + content["start"] + " to " + content["end"]
    dateP.level = 1

    tf.add_paragraph()
    chartNote = tf.add_paragraph()
    chartNote.text = "Chart Value Rounding"
    chartNote.level = 0

    chartNoteP = tf.add_paragraph()
    chartNoteP.text = "Values displayed in the charts are rounded to the nearest visible decimal place."
    chartNoteP.level = 1

    # for paragraph in ph.text_frame.paragraphs:
    #     for run in paragraph.runs:
    #         run.font.size = Pt(12)


def update_divider_slide(the_slide, title):
    update_text_placeholder(the_slide, 1, title)

def update_section_title_slide(the_slide, title, subtitle):
    update_title(the_slide, title)
    update_text_placeholder(the_slide, 1, subtitle)

def update_chart_and_data_slide(the_slide, title, segment, data, chart_path, footer=""):
    update_title(the_slide, title)
    update_text_placeholder(the_slide, 15, f"Segment: {segment}", size=12)
    update_table_placeholder(the_slide, 14, data, size=12)
    update_image_placeholder(the_slide, 13, chart_path)
    update_text_placeholder(the_slide, 16, footer, size=12)

def update_heatmap_slide(the_slide, title, segment, image_path):
    update_title(the_slide, title)
    update_text_placeholder(the_slide, 13, f"{segment}", size=12)
    update_image_placeholder(the_slide, 14, image_path, scale_image = True)

def update_chart_slide(the_slide, title, segment, image_path):
    update_title(the_slide, title)
    update_text_placeholder(the_slide, 13, f"{segment}", size=12)
    update_image_placeholder(the_slide, 14, image_path, scale_image = True)

def generate_report(config):
    if type(config) == type(str):
        config = json.load(open(config))
    Report = create()
    title_slide = get_title_slide(Report)
    update_title_slide(title_slide, f'{config["id"]}', "Post-Test Data Report")
    for content in config["content"]:
        layout_name = content["layout"]
        name = content["name"]
        print(f"Generating {layout_name} slide for {name}")
        layout = Report.slide_layouts.get_by_name(layout_name)
        new_slide = Report.slides.add_slide(layout)

        if layout_name == "Title Slide 1" or layout_name == "Title Slide 2":
            update_title_slide(new_slide, content["title"], content["subtitle"])

        elif layout_name == "Chart and Data":
            update_chart_and_data_slide(new_slide, content["title"], content["segment"], content["data"], content["image_path"], content["footer"])

        elif "Divider Slide 1" in layout_name:
            update_divider_slide(new_slide, content["title"])
            
        elif "Long Form Messaging 1" in layout_name and "Report Info Slide" not in content["type"]:
            update_notes_slide(new_slide, content["title"], content["content"])
            
        elif "Long Form Messaging 1" in layout_name and "Report Info Slide" in content["type"]:
            update_device_report_slide(new_slide, content["title"], content["content"])

        elif "Heatmap" in layout_name:
            update_heatmap_slide(new_slide, content["title"], content["segment"], content["image_path"])

        elif "Chart Only" in layout_name:
            update_chart_slide(new_slide, content["title"], content["segment"], content["image_path"])
    
    save(Report, f'{config["id"]} - Post-Test Data Report')
    utils.print_success("Report Generated Successfully")
